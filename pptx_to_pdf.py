import os
from pptx import Presentation
from fpdf import FPDF

class PDF(FPDF):
    def header(self):
        self.set_font('Arial', 'B', 12)

    def footer(self):
        self.set_y(-15)
        self.set_font('Arial', 'I', 8)
        self.cell(0, 10, f'Page {self.page_no()}', 0, 0, 'C')

def pptx_to_pdf(pptx_path, pdf_path):
    # Load the presentation
    prs = Presentation(pptx_path)
    pdf = PDF()
    pdf.set_auto_page_break(auto=True, margin=15)

    for slide in prs.slides:
        pdf.add_page()
        # Convert slide contents (assuming text) to PDF
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                pdf.set_font('Arial', '', 12)
                pdf.multi_cell(0, 10, shape.text)

    pdf.output(pdf_path)

# Example usage
if __name__ == '__main__':
    pptx_to_pdf('input.pptx', 'output.pdf')
